"""``PptxFormat`` — Microsoft PowerPoint ``.pptx`` (Office Open XML) encoder/decoder.

Reads and writes use ``python-pptx``. ``.pptx`` is a zipped XML bundle:
random access requires the full archive, so this is a
:class:`BatchFileFormat`.

Each slide maps to one record with ``slide_number``, ``text`` (the
concatenation of every text-frame's text on the slide) and optionally
``notes`` (speaker notes, when ``extract_speaker_notes=True``).
Encoding is intentionally simple — every record produces a slide with a
single text frame plus speaker notes.

Round-trip is lossy: ``python-pptx`` cannot recover the exact original
shape geometry from a freshly emitted text frame. Tests should compare
text content (whitespace-normalised) rather than byte equality.

Security: pirn does not sandbox ``python-pptx``. Malformed archives may
trigger upstream library bugs. Treat untrusted payloads accordingly.

Install: ``pip install pirn[pptx]``.
"""

from __future__ import annotations

import io
from collections.abc import Iterable, Mapping
from typing import Any

from pirn.domains.connectors.file_formats.batch_file_format import (
    BatchFileFormat,
)


class PptxFormat(BatchFileFormat):
    """Whole-file PPTX encoder/decoder."""

    def __init__(self, extract_speaker_notes: bool = True) -> None:
        if not isinstance(extract_speaker_notes, bool):
            raise TypeError(
                "PptxFormat: extract_speaker_notes must be a bool, "
                f"got {type(extract_speaker_notes).__name__}"
            )
        self._extract_speaker_notes = extract_speaker_notes

    @property
    def name(self) -> str:
        return "pptx"

    @property
    def extract_speaker_notes(self) -> bool:
        return self._extract_speaker_notes

    async def _decode_full(
        self, payload: bytes
    ) -> Iterable[Mapping[str, Any]]:
        pptx = self._load_pptx()
        presentation = pptx.Presentation(io.BytesIO(payload))
        records: list[Mapping[str, Any]] = []
        for index, slide in enumerate(presentation.slides):
            text = self._slide_text(slide)
            notes: str | None = None
            if self._extract_speaker_notes:
                notes = self._slide_notes(slide)
            records.append(
                {
                    "slide_number": index + 1,
                    "text": text,
                    "notes": notes,
                }
            )
        return records

    async def _encode_full(
        self, records: Iterable[Mapping[str, Any]]
    ) -> bytes:
        pptx = self._load_pptx()
        presentation = pptx.Presentation()
        blank_layout = presentation.slide_layouts[6]
        for record in records:
            if "text" not in record:
                raise ValueError(
                    "PptxFormat: record missing required 'text' field "
                    f"— got keys {list(record.keys())}"
                )
            text = record["text"]
            if not isinstance(text, str):
                raise TypeError(
                    "PptxFormat: 'text' must be a string, got "
                    f"{type(text).__name__}"
                )
            notes = record.get("notes")
            if notes is not None and not isinstance(notes, str):
                raise TypeError(
                    "PptxFormat: 'notes' must be a string or None, "
                    f"got {type(notes).__name__}"
                )
            slide = presentation.slides.add_slide(blank_layout)
            self._add_textbox(presentation, slide, text)
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
        buf = io.BytesIO()
        presentation.save(buf)
        return buf.getvalue()

    @staticmethod
    def _slide_text(slide: Any) -> str:
        parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        parts.append(run.text)
        return "\n".join(parts)

    @staticmethod
    def _slide_notes(slide: Any) -> str | None:
        if not getattr(slide, "has_notes_slide", False):
            return None
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return None
        notes_frame = notes_slide.notes_text_frame
        if notes_frame is None:
            return None
        text = notes_frame.text
        return text if text else None

    @staticmethod
    def _add_textbox(presentation: Any, slide: Any, text: str) -> None:
        from pptx.util import Emu

        # Use the slide's own dimensions (Emu); fall back to a fixed
        # rectangle when geometry is unavailable.
        left = Emu(914400)  # 1 inch
        top = Emu(914400)
        width = presentation.slide_width - Emu(1828800)  # 2 inch margin
        height = presentation.slide_height - Emu(1828800)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = text

    @staticmethod
    def _load_pptx() -> Any:
        try:
            import pptx
        except ImportError as exc:
            raise ImportError(
                "PptxFormat requires python-pptx. Install with "
                "`pip install pirn[pptx]`."
            ) from exc
        return pptx
